import os
from pypdf import PdfReader
from docx import Document
import pandas as pd
from pptx import Presentation
from pptx.shapes.picture import Picture
from services.image_service import analyze_image
import xml.etree.ElementTree as ET
import xmltodict
from bs4 import BeautifulSoup
from striprtf.striprtf import rtf_to_text
import email
from email import policy
from email.parser import BytesParser

def read_file_content(file_path, content_type, llm):

    file_type = ""
    file_content = ""

    if file_path:
        try:
            
            extension = os.path.splitext(file_path)[1].lower()
            
            if extension in [".txt", ".md", ".csv", ".json", ".py", ".js", ".cs"]:
                file_type = "text"
                with open(file_path, "r", encoding="utf-8") as f:
                    file_content = f.read()[:10000]
            
            # PDF
            elif extension == ".pdf":
                file_type = "pdf"
                reader = PdfReader(file_path)
                file_content = ""

                for page in reader.pages:
                    file_content += page.extract_text() or ""
            
            # Word (.docx)
            elif extension in [".doc", ".docx"]:
                file_type = "doc"
                doc = Document(file_path)
                file_content = "\n".join(
                    p.text for p in doc.paragraphs
                )
            
            # Word (.rtf)
            elif extension in [".rtf"]:
                file_type = "rtf"
                with open(file_path, "r", encoding="utf-8", errors="ignore") as file:
                    rtf_content = file.read()

                file_content = rtf_to_text(rtf_content)

            # Excel (.xlsx)
            elif extension in [".xls", ".xlsx"]:
                file_type = "excel"
                df = pd.read_excel(file_path)
                file_content = df.to_string()

            # PPT (.ppt)
            elif extension in [".pptx"]:
                file_type = "ppt"
                prs = Presentation(file_path)

                # # # slides_text = []

                # # # # print("Slides Found:", len(prs.slides))

                # # # for slide_no, slide in enumerate(prs.slides, start=1):
                # # #     # print(f"Slide {slide_no}")

                # # #     slide_content = []

                # # #     for shape in slide.shapes:
                # # #         # print(type(shape))
                # # #         if hasattr(shape, "text"):
                # # #             slide_content.append(shape.text)

                # # #     slides_text.append(
                # # #         f"Slide {slide_no}\n"
                # # #         + "\n".join(slide_content)
                # # #     )

                # # # file_content = "\n\n".join(slides_text)


                slides_data = []
                file_content = ""

                for slide_no, slide in enumerate(prs.slides, start=1):

                    slide_text = []
                    slide_image = []
                    slide_notes = ""

                    for shape_no, shape in slide.shapes:

                        if hasattr(shape, "text"):
                            if shape.text.strip():
                                slide_text.append(shape.text.strip())

                        # print("SHAPE : ")
                        # print(shape)
                        # print("******************************")

                        if hasattr(shape, "image") or shape.shape_type == 13 or isinstance(shape, Picture):  # Picture
                            try:

                                print("In Image Block")

                                image = shape.image

                                # One more useful thing: if your PPT contains many company logos or decorative icons, 
                                # you may want to skip tiny images:
                                # if len(image.blob) < 5000:
                                #     continue

                                image_bytes = image.blob
                                ext = image.ext  # png, jpg, jpeg, gif, etc.

                                image_path = os.path.join("uploads", f"slide_{slide_no}_{shape_no}.{ext}")
                                with open(image_path, "wb") as f:
                                    f.write(image_bytes)

                                imageResult = analyze_image(image_path, llm, "Describe this slide image in detail")
                                slide_image.append(imageResult)

                                print("IMAGE")
                                print(imageResult)
                                print("******************************")

                            except Exception as e:
                                import traceback
                                traceback.print_exc()
                                error_message = repr(e)
                                print("IMAGE ERROR : ")
                                print(error_message)

                    try:
                        if slide.has_notes_slide:
                            slide_notes = slide.notes_slide.notes_text_frame.text
                            
                    except Exception as e:
                        import traceback
                        traceback.print_exc()
                        error_message = repr(e)

                    slides_data.append({
                        "slide_no": slide_no,
                        "text": "\n".join(slide_text),
                        "notes": slide_notes,
                        "IMAGE DESCRIPTION":"\n".join(slide_image)
                    })

                for slide in slides_data:

                    file_content += f"""
                    ===== SLIDE {slide['slide_no']} =====

                    SLIDE CONTENT:
                    {slide['text']}

                    SPEAKER NOTES:
                    {slide['notes']}
                    """

                # print("OK")
                # print(file_content)

            # Image
            elif extension in [".jpg", ".jpeg", ".png", ".webp"]:
                file_type = "image"
                file_content = file_path

            # Audio file
            elif extension in [".mp3", ".wav", ".m4a", ".aac", ".ogg", ".flac"]:
                file_type = "audio"
                file_content = file_path

            # Video file
            elif extension in [".mp4", ".avi", ".mov", ".webm", ".mkv", ".mpeg", ".mpg"] or content_type.startswith("video/"):
                file_type = "video"
                file_content = file_path
            
            # XML file
            elif extension == ".xml":
                file_type = "xml"
                tree = ET.parse(file_path)
                root = tree.getroot()
                output = []

                def walk(node, level=0):
                    indent = "  " * level
                    text = node.text.strip() if node.text else ""
                    output.append(
                        f"{indent}{node.tag}: {text}"
                    )

                    for child in node:
                        walk(child, level + 1)

                walk(root)

                file_content = "XML CONTENT:\n\n" + "\n".join(output)
                
            # HTML file
            elif extension in [".htm", ".html"]:
                file_type = "html"
                
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    html = f.read()
                
                soup = BeautifulSoup(html, "html.parser")

                # Get page Title
                title = ""
                if soup.title:
                    title = soup.title.string

                # Get page LINK
                links = []
                for link in soup.find_all("a"):
                    href = link.get("href")
                    if href:
                        links.append(href)

                # Get page IMAGES
                images = []
                for img in soup.find_all("img"):
                    src = img.get("src")
                    if src:
                        images.append(src)

                text_content = soup.get_text(separator="\n", strip=True)

                file_content = f"""
                    HTML PAGE TITLE:
                    {title}

                    HTML CONTENT:

                    {text_content}
                    """
                
                file_content += "\n\nLINKS:\n"
                file_content += "\n".join(links)

                file_content += "\nIMAGES:\n"
                file_content += "\n".join(images)

            # EMAIL
            elif extension == ".eml":

                file_type = "email"
                with open(file_path, "rb") as f:
                    msg = BytesParser(policy=policy.default).parse(f)
            
                subject = msg.get("subject", "")
                sender = msg.get("from", "")
                recipient = msg.get("to", "")

                body = ""

                if msg.is_multipart():
                    for part in msg.walk():
                        if part.get_content_type() == "text/plain":
                            body += part.get_content()
                else:
                    body = msg.get_content()


                file_content = f"""
                Subject: {subject}
                From: {sender}
                To: {recipient}

                {body}
                """

            # Outlook MSG
            elif extension == ".msg":

                file_type = "email"

                msg = extract_msg.Message(file_path)
                subject = msg.subject
                sender = msg.sender
                body = msg.body

                file_content = f"""
                Subject: {subject}
                From: {sender}

                {body}
                """
            

        except Exception as e:
            file_content = (
                    f"Uploaded file available at: {file_path}\n"
                    f"Could not read content: {str(e)}"
                )


    return {
        "type": file_type,
        "content": file_content
    }